import os
import re
import pickle
import fitz # PyMuPDF
from pptx import Presentation
try:
    from moviepy import VideoFileClip
except ImportError:
    from moviepy.editor import VideoFileClip
import google.generativeai as genai
from dotenv import load_dotenv
import requests
from pinecone import Pinecone, ServerlessSpec
from pinecone_text.sparse import BM25Encoder

# Load environment variables
load_dotenv(dotenv_path="../.env")
import concurrent.futures

# Set up APIs
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
LLM_PROVIDER = os.getenv("LLM_PROVIDER", "gemini")
OLLAMA_BASE_URL = os.getenv("OLLAMA_BASE_URL", "http://localhost:11434")
OLLAMA_EMBED_MODEL = os.getenv("OLLAMA_EMBED_MODEL", "nomic-embed-text")

# Initialize Gemini if API key is present
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)

# Define directories
DATA_DIR = "./data"
DATA_SUBDIRS = {
    "pdf": os.path.join(DATA_DIR, "pdf_data"),
    "pptx": os.path.join(DATA_DIR, "pptx_data"),
    "vtt": os.path.join(DATA_DIR, "vtt_data"),
}
os.makedirs(DATA_DIR, exist_ok=True)
for subdir in DATA_SUBDIRS.values():
    os.makedirs(subdir, exist_ok=True)

# Path to save BM25 encoder model locally
BM25_MODEL_PATH = "./bm25_ollama.json"
DEFAULT_INDEX_NAME = "taxbot-hybrid-index"

class DocumentIngestionPipeline:
    def __init__(self, index_name=DEFAULT_INDEX_NAME):
        self.index_name = index_name
        self.pc = None
        self.index = None
        self.bm25 = BM25Encoder()
        
        if PINECONE_API_KEY:
            self.pc = Pinecone(api_key=PINECONE_API_KEY)

    def init_pinecone_index(self):
        """Initializes or connects to the Pinecone hybrid index."""
        if not self.pc:
            print("[WARNING] Pinecone API Key missing. Skipping index initialization.")
            return False
            
        # For Pinecone hybrid search, the metric MUST be dotproduct
        if self.index_name not in self.pc.list_indexes().names():
            print(f"Creating a new serverless index '{self.index_name}'...")
            self.pc.create_index(
                name=self.index_name,
                dimension=768, # Both nomic-embed-text and text-embedding-004 are 768-dim
                metric="dotproduct",
                spec=ServerlessSpec(
                    cloud="aws",
                    region="us-east-1"
                )
            )
        self.index = self.pc.Index(self.index_name)
        print(f"Connected to Pinecone index: '{self.index_name}'")
        return True

    # --- 1. Document Parsers ---

    def parse_pdf(self, file_path):
        """Extracts text page-by-page using PyMuPDF (fast!)."""
        chunks = []
        filename = os.path.basename(file_path)
        print(f"Parsing PDF: {filename}...")
        
        try:
            doc = fitz.open(file_path)
            for i, page in enumerate(doc):
                text = page.get_text()
                if text and text.strip():
                    page_num = i + 1
                    chunks.append({
                        "text": text,
                        "metadata": {
                            "source": filename,
                            "file_type": "pdf",
                            "page": page_num
                        }
                    })
            doc.close()
        except Exception as e:
            print(f"[ERROR] Failed to parse {filename}: {e}")
            
        return chunks

    def parse_ppt(self, file_path):
        """Extracts text slide-by-slide."""
        chunks = []
        filename = os.path.basename(file_path)
        print(f"Parsing PPT: {filename}...")
        
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            combined_text = "\n".join(slide_text)
            if combined_text.strip():
                slide_num = i + 1
                chunks.append({
                    "text": combined_text,
                    "metadata": {
                        "source": filename,
                        "file_type": "ppt",
                        "slide": slide_num
                    }
                })
        return chunks

    def parse_vtt(self, file_path):
        """Extracts text from WebVTT subtitle files and chunks by timestamp segments."""
        chunks = []
        filename = os.path.basename(file_path)
        print(f"Parsing VTT subtitle: {filename}...")
        
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
                
            # Split blocks by timestamp arrows (e.g. "00:01:02.000 --> 00:01:10.000")
            blocks = re.split(r"(\d{2}:\d{2}:\d{2}\.\d{3}\s*-->\s*\d{2}:\d{2}:\d{2}\.\d{3})", content)
            
            # Parse each timestamped caption block
            captions = []
            for i in range(1, len(blocks), 2):
                timestamp_str = blocks[i]
                if i + 1 < len(blocks):
                    text_block = blocks[i + 1].strip()
                else:
                    continue
                
                # Extract start timestamp and format as [MM:SS] or [HH:MM:SS]
                match = re.match(r"(\d{2}):(\d{2}):(\d{2})", timestamp_str)
                if match:
                    h, m, s = match.groups()
                    timestamp = f"[{h}:{m}:{s}]" if h != "00" else f"[{m}:{s}]"
                else:
                    timestamp = "[00:00]"
                
                # Clean up caption text (remove line numbers, blank lines)
                text_lines = [line.strip() for line in text_block.split("\n") if line.strip() and not line.strip().isdigit()]
                if text_lines:
                    text_content = " ".join(text_lines)
                    captions.append((timestamp, text_content))
            
            # Accumulate captions into ~1000 character paragraph chunks
            accumulated_text = ""
            chunk_timestamp = "[00:00]"
            for ts, txt in captions:
                if not accumulated_text:
                    chunk_timestamp = ts
                    accumulated_text = f"{ts} {txt}"
                else:
                    accumulated_text += f" {txt}"
                    
                if len(accumulated_text) >= 1000:
                    chunks.append({
                        "text": accumulated_text,
                        "metadata": {
                            "source": filename,
                            "file_type": "video",
                            "timestamp": chunk_timestamp
                        }
                    })
                    accumulated_text = ""
                    
            # Don't lose the final accumulated text
            if accumulated_text:
                chunks.append({
                    "text": accumulated_text,
                    "metadata": {
                        "source": filename,
                        "file_type": "video",
                        "timestamp": chunk_timestamp
                    }
                })
                
            print(f"  -> Extracted {len(chunks)} chunks from VTT subtitle.")
        except Exception as e:
            print(f"[ERROR] Failed to parse VTT {filename}: {e}")
            
        return chunks

    def parse_video(self, file_path):
        """Extracts audio, transcribes it via Gemini Files API, and chunks the transcript."""
        chunks = []
        filename = os.path.basename(file_path)
        print(f"Processing Video: {filename}...")
        
        # 1. Extract audio file
        audio_path = os.path.splitext(file_path)[0] + ".mp3"
        try:
            video = VideoFileClip(file_path)
            if video.audio:
                print(f"Extracting audio to {audio_path}...")
                video.audio.write_audiofile(audio_path, logger=None)
            else:
                print("[WARNING] Video file has no audio track.")
                return chunks
        except Exception as e:
            print(f"[WARNING] Failed to extract audio: {e}")
            return chunks
        
        # 2. Upload and transcribe via Gemini
        if not GEMINI_API_KEY:
            print("[WARNING] Gemini API Key missing. Video transcription requires Gemini Files API.")
            return chunks
            
        try:
            print(f"Uploading audio to Gemini Files API...")
            audio_file = genai.upload_file(path=audio_path)
            
            # Wait for file processing if needed
            print("Transcribing audio file via Gemini...")
            model = genai.GenerativeModel("gemini-3.1-flash-lite")
            response = model.generate_content([
                "Please transcribe this tax lecture audio file accurately. "
                "Output timestamps in brackets like [MM:SS] at the start of each logical paragraph, "
                "so we can cross-reference the discussion point back to the video time.",
                audio_file
            ])
            
            transcript = response.text
            print("Transcribed successfully! Parsing timestamps and chunking...")
            
            # Save the transcript for offline viewing
            transcript_txt_path = os.path.splitext(file_path)[0] + "_transcript.txt"
            with open(transcript_txt_path, "w", encoding="utf-8") as f:
                f.write(transcript)
            
            # Clean up the audio file on Gemini side
            genai.delete_file(audio_file.name)
            
            # Remove local audio copy
            if os.path.exists(audio_path):
                os.remove(audio_path)

            # Chunk the transcript by timestamp segments
            # Find all timestamps like [01:23] or [12:34]
            segments = re.split(r"(\[\d{2}:\d{2}\])", transcript)
            
            current_timestamp = "[00:00]"
            for part in segments:
                if re.match(r"^\[\d{2}:\d{2}\]$", part):
                    current_timestamp = part
                else:
                    text_content = part.strip()
                    if text_content:
                        chunks.append({
                            "text": f"{current_timestamp} {text_content}",
                            "metadata": {
                                "source": filename,
                                "file_type": "video",
                                "timestamp": current_timestamp
                            }
                        })
                        
        except Exception as e:
            print(f"[WARNING] Error during video transcription: {e}")
            
        return chunks

    # --- 2. Chunk Splitting (Refining chunks for embeddings) ---

    def split_into_embeddings_chunks(self, raw_chunks, max_chars=1000, overlap=200):
        """Splits raw page/slide text into smaller overlapping chunks if they exceed size limits."""
        refined_chunks = []
        for raw in raw_chunks:
            text = raw["text"]
            metadata = raw["metadata"]
            
            if len(text) <= max_chars:
                refined_chunks.append(raw)
                continue
                
            # Slit text into overlapping windows
            start = 0
            chunk_idx = 1
            while start < len(text):
                end = start + max_chars
                segment = text[start:end]
                
                # Copy metadata and add chunk index
                chunk_meta = metadata.copy()
                chunk_meta["chunk_index"] = chunk_idx
                
                refined_chunks.append({
                    "text": segment,
                    "metadata": chunk_meta
                })
                
                start += (max_chars - overlap)
                chunk_idx += 1
                
        return refined_chunks

    # --- 3. Dense Embeddings ---

    def get_dense_embedding(self, text):
        """Generates dense vector embeddings using Gemini or local Ollama."""
        if LLM_PROVIDER == "ollama":
            try:
                res = requests.post(f"{OLLAMA_BASE_URL}/api/embeddings", json={
                    "model": OLLAMA_EMBED_MODEL,
                    "prompt": text
                })
                res.raise_for_status()
                return res.json()["embedding"]
            except Exception as e:
                raise RuntimeError(f"Ollama embedding failed: {e}. Ensure Ollama is running and has pulled {OLLAMA_EMBED_MODEL}")
        else:
            if not GEMINI_API_KEY:
                raise ValueError("Gemini API key is not configured.")
            res = genai.embed_content(
                model="models/text-embedding-004",
                content=text,
                task_type="retrieval_document",
                output_dimensionality=768
            )
            return res["embedding"]

    # --- 4. Sparse Vector Fitting & Encoding ---

    def fit_bm25(self, all_texts):
        """Fits the BM25 model on all text chunks and saves it locally."""
        print("Fitting BM25 encoder on corpus...")
        # To avoid Out-Of-Memory (OOM) on 512MB RAM servers, fit on a subset if corpus is huge
        sample_texts = all_texts[:5000] if len(all_texts) > 5000 else all_texts
        self.bm25.fit(sample_texts)
        self.bm25.dump(BM25_MODEL_PATH)
        print(f"BM25 encoder fitted and saved to: {BM25_MODEL_PATH}")

    def load_bm25(self):
        """Loads a pre-fitted local BM25 model."""
        if os.path.exists(BM25_MODEL_PATH):
            self.bm25.load(BM25_MODEL_PATH)
            return True
        return False

    def get_sparse_embedding(self, text):
        """Generates sparse vector representations using the BM25 model."""
        # Returns a dict of indices and values required by Pinecone
        return self.bm25.encode_documents(text)

    # --- 5. Main Processing Run ---

    def process_all_files(self):
        """Iterates over data directory (including subdirs), processes files, and pushes to Pinecone."""
        if not os.path.exists(DATA_DIR):
            print(f"Data directory '{DATA_DIR}' does not exist.")
            return

        all_raw_chunks = []
        
        # Collect all files from subdirectories and root data dir (backwards compatibility)
        all_file_paths = []
        scan_dirs = [DATA_DIR] + list(DATA_SUBDIRS.values())
        for scan_dir in scan_dirs:
            if not os.path.exists(scan_dir):
                continue
            for filename in os.listdir(scan_dir):
                file_path = os.path.join(scan_dir, filename)
                if os.path.isfile(file_path):
                    all_file_paths.append(file_path)
        
        # Build a set of recording base names that have a VTT subtitle file
        # e.g. "GMT20260108-120125_Recording.cc.vtt" -> base "GMT20260108-120125_Recording"
        vtt_bases = set()
        for file_path in all_file_paths:
            if file_path.lower().endswith(".vtt"):
                base = os.path.basename(file_path).split(".")[0]
                vtt_bases.add(base)
        
        for file_path in all_file_paths:
            filename = os.path.basename(file_path)
            ext = os.path.splitext(filename)[1].lower()
            
            try:
                if ext == ".pdf":
                    all_raw_chunks.extend(self.parse_pdf(file_path))
                elif ext in [".ppt", ".pptx"]:
                    all_raw_chunks.extend(self.parse_ppt(file_path))
                elif ext == ".vtt":
                    # Parse VTT subtitle files directly (instant, free, accurate)
                    all_raw_chunks.extend(self.parse_vtt(file_path))
                elif ext in [".mp4", ".avi", ".mkv", ".mov", ".m4a", ".mp3", ".wav"]:
                    # De-duplicate: skip media files if a VTT subtitle exists for this recording
                    base = filename.split(".")[0]
                    # Also strip common suffixes like "_2880x1800" or "_gallery_2880x1800"
                    stripped_base = re.sub(r'(_gallery)?_\d+x\d+$', '', base)
                    if base in vtt_bases or stripped_base in vtt_bases:
                        print(f"Skipping '{filename}' (VTT subtitle already provides this transcript)")
                        continue
                    # No VTT available — fall back to Gemini-based video transcription
                    all_raw_chunks.extend(self.parse_video(file_path))
                elif ext == ".png":
                    continue  # Silently skip image files
                else:
                    print(f"Ignoring unsupported file type: {filename}")
            except Exception as parse_err:
                print(f"[WARNING] Failed to parse file '{filename}': {parse_err}. Skipping...")

        if not all_raw_chunks:
            print("No valid files processed or data folder is empty.")
            return

        # Refine text into smaller embedding chunks
        refined_chunks = self.split_into_embeddings_chunks(all_raw_chunks)
        print(f"Generated {len(refined_chunks)} total chunks.")

        # Fit BM25 on the corpus texts
        corpus_texts = [c["text"] for c in refined_chunks]
        self.fit_bm25(corpus_texts)

        # Initialize Pinecone
        if not self.init_pinecone_index():
            print("[ERROR] Cannot proceed with upload: Pinecone connection not initialized.")
            return

        # Generate vectors and upsert to Pinecone
        print("Generating embeddings and uploading to Pinecone in batches...")
        import time
        vectors_to_upsert = []
        
        # Filter out empty chunks first
        valid_chunks = [(idx, chunk) for idx, chunk in enumerate(refined_chunks) if chunk["text"].strip()]
        
        batch_size = 90  # Safe limit under 100 for Gemini API
        
        for i in range(0, len(valid_chunks), batch_size):
            batch = valid_chunks[i:i+batch_size]
            texts = [c["text"] for idx, c in batch]
            
            try:
                # 1. Dense Embeddings
                if LLM_PROVIDER == "gemini":
                    if not GEMINI_API_KEY:
                        raise ValueError("Gemini API key is not configured.")
                    res = genai.embed_content(
                        model="models/text-embedding-004",
                        content=texts,
                        task_type="retrieval_document",
                        output_dimensionality=768
                    )
                    dense_vectors = res["embedding"]
                else:
                    with concurrent.futures.ThreadPoolExecutor(max_workers=20) as executor:
                        dense_vectors = list(executor.map(self.get_dense_embedding, texts))
                
                # 2. Prepare Pinecone Vectors
                for j, (idx, chunk) in enumerate(batch):
                    text = chunk["text"]
                    metadata = chunk["metadata"]
                    metadata["text"] = text
                    
                    sparse_vector = self.get_sparse_embedding(text)
                    chunk_id = f"chunk_{metadata['source']}_{idx}"
                    chunk_id = re.sub(r'[^a-zA-Z0-9_\-\.#]', '_', chunk_id)
                    
                    vector_data = {
                        "id": chunk_id,
                        "values": dense_vectors[j],
                        "metadata": metadata
                    }
                    if sparse_vector.get("indices") and sparse_vector.get("values"):
                        vector_data["sparse_values"] = sparse_vector
                        
                    vectors_to_upsert.append(vector_data)
                
                # 3. Upsert
                if vectors_to_upsert:
                    self.index.upsert(vectors=vectors_to_upsert)
                    print(f"Uploaded batch of {len(vectors_to_upsert)} chunks... ({i+len(batch)}/{len(valid_chunks)})")
                    vectors_to_upsert = []
                    
            except Exception as e:
                print(f"[ERROR] Failed to process batch {i}: {e}")
                
        print("[SUCCESS] Ingestion Pipeline completed successfully!")
        
        # Return refined_chunks so callers can use it for tracking
        return refined_chunks

if __name__ == "__main__":
    pipeline = DocumentIngestionPipeline()
    refined_chunks = pipeline.process_all_files()
    
    # Track ingested documents in Supabase so they appear in the Admin Panel
    if refined_chunks:
        try:
            from database import upsert_ingested_document
            
            # Count chunks per source file
            chunks_per_file = {}
            for chunk in refined_chunks:
                source = chunk["metadata"].get("source", "unknown")
                chunks_per_file[source] = chunks_per_file.get(source, 0) + 1
            
            # Upsert each file's metadata into Supabase
            for filename, chunk_count in chunks_per_file.items():
                file_path = os.path.join(DATA_DIR, filename)
                ext = os.path.splitext(filename)[1].lower().replace(".", "").upper()
                size_mb = os.path.getsize(file_path) / (1024 * 1024) if os.path.exists(file_path) else 0
                upsert_ingested_document(filename, ext, size_mb, chunk_count)
            
            print(f"[SUCCESS] Tracked {len(chunks_per_file)} documents in Supabase.")
        except Exception as e:
            print(f"[WARNING] Failed to track documents in Supabase: {e}")
